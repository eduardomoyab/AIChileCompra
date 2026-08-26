import os
import io
import csv
import logging
import shutil
import unicodedata
import zipfile
import gc
import numpy as np
from tqdm import tqdm
from docx import Document
from openpyxl import load_workbook
from bs4 import BeautifulSoup
from typing import Optional, List
import fitz  # PyMuPDF
import easyocr
from PIL import Image

try:
    import rarfile
    _rarfile_available = True
except ImportError:
    _rarfile_available = False
    logging.warning("rarfile no disponible — archivos .rar no serán extraídos")


def _find_unrar_tool() -> Optional[str]:
    """Busca el ejecutable unrar/UnRAR en orden de prioridad."""
    # 1. Variable de entorno explícita
    env_path = os.environ.get("UNRAR_PATH", "")
    if env_path and os.path.isfile(env_path):
        return env_path
    # 2. Ubicaciones típicas de WinRAR en Windows
    for candidate in [
        r"C:\Program Files\WinRAR\UnRAR.exe",
        r"C:\Program Files (x86)\WinRAR\UnRAR.exe",
    ]:
        if os.path.isfile(candidate):
            return candidate
    # 3. PATH del sistema (Linux: apt-get install unrar)
    for name in ("unrar", "unrar-free"):
        found = shutil.which(name)
        if found:
            return found
    return None

try:
    import xlrd
    _xlrd_available = True
except ImportError:
    _xlrd_available = False
    logging.warning("xlrd no disponible — archivos .xls no serán procesados")

try:
    from pptx import Presentation
    _pptx_available = True
except ImportError:
    _pptx_available = False
    logging.warning("python-pptx no disponible — archivos .pptx no serán procesados")

# Configuración por defecto
MAX_PAGES = 15
TEXT_THRESHOLD = 100
MAX_FILENAME_LENGTH = 100

# Reader EasyOCR compartido entre todos los AttachmentProcessor del proceso
# Se inicializa una vez al importar el módulo (una vez por worker de gunicorn)
try:
    import torch
    _gpu = torch.cuda.is_available()
except ImportError:
    _gpu = False

try:
    _easyocr_reader = easyocr.Reader(['es', 'en'], gpu=_gpu)
    logging.info(f"EasyOCR inicializado al cargar módulo (gpu={_gpu})")
except Exception as _e:
    _easyocr_reader = None
    logging.warning(f"EasyOCR no disponible: {_e}")


class AttachmentProcessor:
    """Procesa archivos adjuntos extrayendo texto de diversos formatos.
    PDFs e imágenes: PyMuPDF (texto nativo) + EasyOCR (OCR para páginas escaneadas).
    """

    def __init__(self, attachments_path: str, output_path: Optional[str] = None,
                 blacklist: Optional[List[str]] = None, use_gpu: bool = False):
        """
        Inicializa el procesador de adjuntos.

        Args:
            attachments_path (str): Ruta a la carpeta con los adjuntos descargados
            output_path (str, optional): Ruta donde guardar los archivos procesados
            blacklist (list, optional): Lista de palabras para filtrar archivos
            use_gpu (bool): Si usar GPU para EasyOCR
        """
        self.attachments_path = os.path.abspath(attachments_path)

        if output_path is None:
            parent_dir = os.path.dirname(self.attachments_path)
            self.output_path = os.path.join(parent_dir, 'processed')
        else:
            self.output_path = os.path.abspath(output_path)

        self.blacklist = blacklist if blacklist else []

        # Reutilizar el reader global inicializado al cargar el módulo
        self.reader = _easyocr_reader

        # Archivos omitidos
        self.skipped_files = set()
        self.skip_files_path = os.path.join(self.output_path, 'skipped_files.txt')

        # Crear directorio de salida
        os.makedirs(self.output_path, exist_ok=True)

        # Cargar archivos previamente omitidos
        self._load_skipped_files()

    def _extract_archives_in_path(self) -> dict:
        """
        Extrae archivos .rar y .zip encontrados en attachments_path al mismo directorio
        (extracción plana, sin estructura de subdirectorios) para que el loop principal
        los procese como si fueran archivos normales.

        Returns:
            dict con: archivos_comprimidos_encontrados, extraidos, fallidos (lista de nombres)
        """
        unrar_tool = _find_unrar_tool()
        if _rarfile_available and unrar_tool:
            rarfile.UNRAR_TOOL = unrar_tool

        found, extracted, failed = [], 0, []

        try:
            entries = os.listdir(self.attachments_path)
        except Exception:
            return {"archivos_comprimidos_encontrados": 0, "extraidos": 0, "fallidos": []}

        for fname in entries:
            fpath = os.path.join(self.attachments_path, fname)
            lower = fname.lower()

            if lower.endswith(".rar"):
                found.append(fname)
                if not _rarfile_available or not unrar_tool:
                    logging.warning(f"RAR encontrado pero unrar no disponible: {fname}")
                    failed.append(fname)
                    continue
                try:
                    rf = rarfile.RarFile(fpath)
                    for member in rf.infolist():
                        if member.is_dir():
                            continue
                        member_name = os.path.basename(member.filename)
                        if not member_name:
                            continue
                        dest = os.path.join(self.attachments_path, member_name)
                        if not os.path.exists(dest):
                            data = rf.read(member)
                            with open(dest, "wb") as out:
                                out.write(data)
                    extracted += 1
                    logging.info(f"✓ RAR extraído: {fname}")
                except Exception as e:
                    logging.warning(f"Error extrayendo RAR {fname}: {e}")
                    failed.append(fname)

            elif lower.endswith(".zip"):
                found.append(fname)
                try:
                    with zipfile.ZipFile(fpath, "r") as zf:
                        for member in zf.infolist():
                            if member.filename.endswith("/"):
                                continue
                            member_name = os.path.basename(member.filename)
                            if not member_name:
                                continue
                            dest = os.path.join(self.attachments_path, member_name)
                            if not os.path.exists(dest):
                                data = zf.read(member)
                                with open(dest, "wb") as out:
                                    out.write(data)
                    extracted += 1
                    logging.info(f"✓ ZIP extraído: {fname}")
                except Exception as e:
                    logging.warning(f"Error extrayendo ZIP {fname}: {e}")
                    failed.append(fname)

        if found:
            logging.info(f"Archivos comprimidos: {len(found)} encontrados, {extracted} extraídos, {len(failed)} fallidos")

        return {
            "archivos_comprimidos_encontrados": len(found),
            "extraidos": extracted,
            "fallidos": failed,
        }

    def _load_skipped_files(self):
        """Carga la lista de archivos previamente omitidos"""
        if os.path.exists(self.skip_files_path):
            with open(self.skip_files_path, 'r', encoding='utf-8') as file:
                self.skipped_files = set(line.strip() for line in file)

    def _add_to_skipped_files(self, file_name: str):
        """Agrega un archivo a la lista de omitidos"""
        if file_name not in self.skipped_files:
            with open(self.skip_files_path, 'a', encoding='utf-8') as file:
                file.write(f"{file_name}\n")
            self.skipped_files.add(file_name)

    def _normalize_filename(self, filename: str) -> str:
        """Normaliza el nombre del archivo removiendo caracteres especiales"""
        normalized = unicodedata.normalize('NFKD', filename).encode('ASCII', 'ignore').decode('ASCII')
        normalized = normalized.replace(' ', '_')

        if len(normalized) > MAX_FILENAME_LENGTH:
            normalized = normalized[:MAX_FILENAME_LENGTH]
            logging.warning(f"Filename truncated to {MAX_FILENAME_LENGTH} chars: {normalized}")

        return normalized

    def _contains_blacklisted_word(self, text: str) -> bool:
        """Verifica si el texto contiene palabras de la blacklist"""
        if not self.blacklist or not text:
            return False
        return any(word.lower() in text.lower() for word in self.blacklist)

    def _ocr_image(self, img: Image.Image) -> str:
        """Aplica EasyOCR a una imagen PIL y retorna el texto."""
        if not self.reader:
            return ''
        try:
            result = self.reader.readtext(np.array(img), detail=0)
            return '\n'.join(result)
        except Exception as e:
            logging.warning(f"EasyOCR error: {e}")
            return ''

    def _extract_text_from_pdf(self, pdf_path: str) -> Optional[str]:
        """
        Extrae texto de un PDF página por página con PyMuPDF.
        - Páginas con texto nativo suficiente → extracción directa.
        - Páginas escaneadas (poco texto)     → render a imagen + EasyOCR.
        """
        try:
            doc = fitz.open(pdf_path)
            pages_text = []

            for i, page in enumerate(doc):
                if i >= MAX_PAGES:
                    break
                native = page.get_text() or ''

                if len(native.strip()) >= TEXT_THRESHOLD:
                    pages_text.append(native)
                else:
                    # Página escaneada: render a 200 DPI y OCR
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img = Image.open(io.BytesIO(pix.tobytes("png")))
                    ocr_text = self._ocr_image(img)
                    pages_text.append(ocr_text)

            doc.close()
            combined = '\n\n'.join(t for t in pages_text if t.strip())
            return combined if combined.strip() else None

        except Exception as e:
            logging.error(f"Error reading PDF {pdf_path}: {e}")
            return None

    def _extract_text_with_ocr(self, image_path: str) -> Optional[str]:
        """Extrae texto de imágenes usando EasyOCR"""
        try:
            img = Image.open(image_path)
            text = self._ocr_image(img)
            return text if text.strip() else None
        except Exception as e:
            logging.error(f"OCR error for {image_path}: {e}")
            return None

    def _extract_text_from_docx(self, docx_path: str) -> Optional[str]:
        """Extrae texto de archivos DOCX (texto nativo sin OCR)"""
        try:
            doc = Document(docx_path)
            text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            return text if text.strip() else None
        except Exception as e:
            logging.error(f"Error extracting text from DOCX {docx_path}: {e}")
            return None

    def _extract_text_from_html(self, html_path: str) -> Optional[str]:
        """Extrae texto de archivos HTML"""
        try:
            with open(html_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file, 'html.parser')
                text = soup.get_text(separator='\n')
            return text if text.strip() else None
        except Exception as e:
            logging.error(f"Error extracting text from HTML {html_path}: {e}")
            return None

    def _extract_text_from_xlsx(self, xlsx_path: str) -> Optional[str]:
        """Extrae texto de archivos Excel"""
        try:
            workbook = load_workbook(xlsx_path, data_only=True)
            text = ""
            for sheet in workbook.sheetnames:
                worksheet = workbook[sheet]
                for row in worksheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
            return text if text.strip() else None
        except Exception as e:
            logging.error(f"Error extracting text from XLSX {xlsx_path}: {e}")
            return None

    def _extract_text_from_xls(self, xls_path: str) -> Optional[str]:
        """Extrae texto de archivos Excel antiguos (.xls)"""
        if not _xlrd_available:
            return None
        try:
            workbook = xlrd.open_workbook(xls_path)
            text = ""
            for sheet in workbook.sheets():
                for row_idx in range(sheet.nrows):
                    row_text = "\t".join([str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)])
                    text += row_text + "\n"
            return text if text.strip() else None
        except Exception as e:
            logging.error(f"Error extracting text from XLS {xls_path}: {e}")
            return None

    def _extract_text_from_pptx(self, pptx_path: str) -> Optional[str]:
        """Extrae texto de presentaciones PowerPoint (.pptx)"""
        if not _pptx_available:
            return None
        try:
            prs = Presentation(pptx_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text += " ".join([run.text for run in para.runs]) + "\n"
            return text if text.strip() else None
        except Exception as e:
            logging.error(f"Error extracting text from PPTX {pptx_path}: {e}")
            return None

    def _extract_text_from_txt(self, txt_path: str) -> Optional[str]:
        """Lee archivos de texto plano"""
        for encoding in ('utf-8', 'latin-1', 'cp1252'):
            try:
                with open(txt_path, 'r', encoding=encoding) as f:
                    text = f.read()
                return text if text.strip() else None
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logging.error(f"Error reading TXT {txt_path}: {e}")
                return None
        return None

    def _extract_text_from_csv(self, csv_path: str) -> Optional[str]:
        """Lee archivos CSV como texto tabular"""
        for encoding in ('utf-8', 'latin-1', 'cp1252'):
            try:
                with open(csv_path, 'r', encoding=encoding, newline='') as f:
                    reader = csv.reader(f)
                    rows = ["\t".join(row) for row in reader]
                return "\n".join(rows) if rows else None
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logging.error(f"Error reading CSV {csv_path}: {e}")
                return None
        return None

    def _extract_text_from_file(self, file_path: str) -> Optional[str]:
        """Extrae texto según el tipo de archivo"""
        ext = file_path.lower()

        if ext.endswith('.pdf'):
            return self._extract_text_from_pdf(file_path)
        elif ext.endswith(('.jpg', '.jpeg', '.png')):
            return self._extract_text_with_ocr(file_path)
        elif ext.endswith('.docx'):
            return self._extract_text_from_docx(file_path)
        elif ext.endswith('.html'):
            return self._extract_text_from_html(file_path)
        elif ext.endswith('.xlsx'):
            return self._extract_text_from_xlsx(file_path)
        elif ext.endswith('.xls'):
            return self._extract_text_from_xls(file_path)
        elif ext.endswith('.pptx'):
            return self._extract_text_from_pptx(file_path)
        elif ext.endswith('.txt'):
            return self._extract_text_from_txt(file_path)
        elif ext.endswith('.csv'):
            return self._extract_text_from_csv(file_path)

        return None

    def _save_text_file(self, txt_path: str, codigo_cotizacion: str,
                       rut_proveedor: str, original_filename: str, content: str):
        """Guarda el contenido extraído en un archivo de texto"""
        with open(txt_path, 'w', encoding='utf-8') as txt_file:
            txt_file.write(f"Codigo Cotizacion: {codigo_cotizacion}\n")
            txt_file.write(f"Rut Proveedor: {rut_proveedor}\n")
            txt_file.write(f"Nombre Original: {original_filename}\n\n")
            txt_file.write("Contenido:\n")
            txt_file.write(content)
        logging.info(f"Text file saved: {txt_path}")

    def process_attachments(self) -> dict:
        """
        Procesa todos los archivos en la carpeta de adjuntos.

        Returns:
            dict: Resultado del procesamiento con:
                - success (bool): Si el proceso fue exitoso
                - processed_files (int): Cantidad de archivos procesados
                - skipped_files (int): Cantidad de archivos omitidos
                - output_path (str): Ruta donde se guardaron los archivos procesados
                - error (str, optional): Mensaje de error si falló
        """
        if not os.path.exists(self.attachments_path):
            return {
                "success": False,
                "error": f"Path does not exist: {self.attachments_path}",
                "processed_files": 0,
                "skipped_files": 0,
                "output_path": self.output_path
            }

        processed_count = 0
        skipped_count = 0
        errors = []

        # Extraer archivos comprimidos antes de procesar
        archive_stats = self._extract_archives_in_path()

        try:
            files_to_process = []

            for file in os.listdir(self.attachments_path):
                file_path = os.path.join(self.attachments_path, file)

                if not file.lower().endswith(('.pdf', '.jpg', '.jpeg', '.png', '.docx', '.html', '.xlsx', '.xls', '.pptx', '.txt', '.csv')):
                    continue

                base_name = os.path.splitext(file)[0]
                normalized_name = self._normalize_filename(base_name)
                txt_path = os.path.join(self.output_path, f"{normalized_name}.txt")

                if os.path.exists(txt_path):
                    logging.info(f"Already processed: {file}")
                    continue

                if file in self.skipped_files:
                    logging.info(f"Previously skipped: {file}")
                    skipped_count += 1
                    continue

                files_to_process.append((file, file_path, txt_path))

            for file, file_path, txt_path in tqdm(files_to_process, desc="Processing attachments"):
                try:
                    text_content = self._extract_text_from_file(file_path)

                    if self.blacklist and text_content and self._contains_blacklisted_word(text_content):
                        logging.info(f"Skipped (blacklist): {file}")
                        self._add_to_skipped_files(file)
                        skipped_count += 1
                        continue

                    if text_content is None or len(text_content.strip()) == 0:
                        logging.warning(f"No text extracted: {file}")
                        self._add_to_skipped_files(file)
                        skipped_count += 1
                        continue

                    # Formato esperado: attachments/{codigo_cotizacion}/{rut_proveedor}/archivo.ext
                    path_parts = self.attachments_path.split(os.sep)
                    if len(path_parts) >= 2:
                        rut_proveedor = path_parts[-1]
                        codigo_cotizacion = path_parts[-2]
                    else:
                        rut_proveedor = "unknown"
                        codigo_cotizacion = "unknown"

                    self._save_text_file(txt_path, codigo_cotizacion, rut_proveedor, file, text_content)
                    processed_count += 1

                    gc.collect()

                except Exception as e:
                    error_msg = f"Error processing {file}: {e}"
                    logging.error(error_msg)
                    errors.append(error_msg)
                    skipped_count += 1

            return {
                "success": True,
                "processed_files": processed_count,
                "skipped_files": skipped_count,
                "output_path": self.output_path,
                "errors": errors if errors else None,
                "archivos_comprimidos": archive_stats,
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "processed_files": processed_count,
                "skipped_files": skipped_count,
                "output_path": self.output_path
            }


def process_attachments_simple(attachments_path: str, output_path: Optional[str] = None,
                               blacklist: Optional[List[str]] = None, use_gpu: bool = False) -> dict:
    """
    Función simple para procesar adjuntos sin necesidad de gestionar la clase.

    Args:
        attachments_path (str): Ruta a la carpeta con los adjuntos descargados
        output_path (str, optional): Ruta donde guardar archivos procesados
        blacklist (list, optional): Lista de palabras para filtrar
        use_gpu (bool): Ignorado (mantenido por compatibilidad)

    Returns:
        dict: Resultado del procesamiento

    Example:
        >>> result = process_attachments_simple("attachments/12345678/76123456-7")
        >>> if result['success']:
        >>>     print(f"Processed {result['processed_files']} files")
        >>>     print(f"Output: {result['output_path']}")
    """
    processor = AttachmentProcessor(attachments_path, output_path, blacklist, use_gpu)
    return processor.process_attachments()
